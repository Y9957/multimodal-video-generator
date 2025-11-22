"""
ppt_parser.py

- PPTX 파일에서 슬라이드별 텍스트 / 표 / 이미지 추출
- 각 슬라이드를 PNG 이미지로 렌더링
- 결과를 SlideData 리스트로 state["slides"]에 담아주는 모듈

다른 모듈에서는 node_parse_ppt(state)를 호출해서 사용하면 됩니다.
"""

import os
import re
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import List, Dict, Optional, TypedDict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ------------------------------------------------------------
# 유틸 함수
# ------------------------------------------------------------

def clean_text(s: str) -> str:
    """여러 줄/공백으로 섞여 있는 텍스트를 한 줄로 정리."""
    return re.sub(r"\s+", " ", s).strip()


def export_slide_as_png(state: Dict, idx: int, dpi: int = 220) -> Dict:
    """
    PPTX를 슬라이드 단위 PNG로 변환하여 state["slide_image"]에 경로를 저장.
    - LibreOffice + pdftoppm 사용 (Colab / Linux 환경 기준)
    - state["work_dir"], state["pptx_path"]가 세팅되어 있어야 함.
    """
    work_dir = Path(state["work_dir"]).expanduser().resolve()
    work_dir.mkdir(parents=True, exist_ok=True)

    pptx = Path(state["pptx_path"]).expanduser().resolve()
    if not pptx.exists():
        raise FileNotFoundError(f"PPTX 없음: {pptx}")

    page_no = idx + 1
    out_prefix = work_dir / "slide_img"

    # 폰트/로케일 설정
    env = os.environ.copy()
    env.update({
        "LANG": "ko_KR.UTF-8",
        "LC_ALL": "ko_KR.UTF-8",
    })

    # ---------- A) 직접 PNG 변환 시도 ----------
    before_png = set(work_dir.glob("*.png"))
    png_cmd = [
        "soffice", "--headless",
        "-env:UserInstallation=file:///tmp/lo_profile",
        "--convert-to", "png:impress_png_Export",
        "--outdir", str(work_dir),
        str(pptx),
    ]
    res_png = subprocess.run(png_cmd, capture_output=True, text=True, env=env)

    created_png = [p for p in work_dir.glob("*.png") if p not in before_png]

    candidate = None
    exact = [p for p in created_png if p.stem.endswith(f"-{page_no}")]
    if exact:
        candidate = max(exact, key=lambda p: p.stat().st_mtime)
    elif created_png:
        candidate = max(created_png, key=lambda p: p.stat().st_mtime)

    if candidate and candidate.exists():
        state["slide_image"] = str(candidate)
        return state

    # ---------- B) 폴백: PPTX → PDF → PNG ----------
    target_pdf = work_dir / f"{pptx.stem}.pdf"
    before_pdf = {p.name for p in work_dir.glob("*.pdf")}

    lo_cmd = [
        "soffice", "--headless",
        "-env:UserInstallation=file:///tmp/lo_profile",
        "--convert-to", "pdf:impress_pdf_Export",
        "--outdir", str(work_dir),
        str(pptx),
    ]
    res_pdf = subprocess.run(lo_cmd, capture_output=True, text=True, env=env)

    if target_pdf.exists():
        pdf_path = target_pdf
    else:
        created = [p for p in work_dir.glob("*.pdf") if p.name not in before_pdf]
        if created:
            pdf_path = max(created, key=lambda p: p.stat().st_mtime)
        else:
            print("LibreOffice 변환 실패")
            print("stdout:", res_pdf.stdout)
            print("stderr:", res_pdf.stderr)
            raise RuntimeError("PPTX → PDF 변환 실패")

    ppm_cmd = [
        "pdftoppm",
        "-f", str(page_no), "-l", str(page_no),
        "-png", "-r", str(dpi),
        str(pdf_path),
        str(out_prefix),
    ]
    subprocess.run(ppm_cmd, capture_output=True, text=True, env=env)

    png_path = Path(f"{out_prefix}-{page_no}.png")
    state["slide_image"] = str(png_path)
    return state


# ------------------------------------------------------------
# 데이터 구조 정의
# ------------------------------------------------------------

@dataclass
class SlideData:
    """한 슬라이드에 대한 구조화된 정보."""
    page: int                          # 페이지 번호 (0-based)
    slide_image: str                   # 슬라이드 스냅샷 PNG 경로
    texts: List[str]                   # 슬라이드 내 텍스트들
    images: List[str]                  # 추출된 이미지 파일 경로들
    tables: List[List[List[str]]]      # 표 데이터 (테이블 리스트 → 행 리스트 → 셀 텍스트)
    summary: Optional[str] = None      # LLM 요약문
    script: Optional[str] = None       # 강의 스크립트
    script_file: Optional[str] = None  # 스크립트 저장 파일 경로
    audio: Optional[str] = None        # TTS 음성 파일 경로
    video: Optional[str] = None        # 슬라이드별 영상 파일 경로


class State(TypedDict, total=False):
    """전체 파이프라인에서 사용하는 최소 State 정의 (이 모듈에서 쓰는 부분만)."""
    pptx_path: str
    work_dir: str
    media_dir: str
    slides: List[SlideData]
    slide_image: str  # export_slide_as_png에서 임시 사용


# ------------------------------------------------------------
# 메인 노드 함수
# ------------------------------------------------------------

def node_parse_ppt(state: State) -> State:
    """
    PPT에서 텍스트, 표, 이미지를 페이지별로 추출하고
    각 슬라이드 PNG도 생성하여 SlideData 리스트로 state["slides"]에 저장한다.

    필요한 state 키:
        - pptx_path: PPTX 파일 경로
        - work_dir:  작업용 디렉토리
        - media_dir: 이미지 등 저장할 디렉토리
    """
    prs = Presentation(state["pptx_path"])
    slides_data: List[SlideData] = []

    for i, slide in enumerate(prs.slides):
        texts: List[str] = []
        tables: List[List[List[str]]] = []
        images: List[str] = []

        for shape in slide.shapes:
            # 텍스트 추출
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = clean_text(paragraph.text)
                    if t:
                        texts.append(t)

            # 표 추출
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data: List[List[str]] = []
                for row in shape.table.rows:
                    row_data = [clean_text(cell.text) for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)

            # 이미지 추출
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext or "png"
                image_filename = f"{state['media_dir']}/slide{i}_img{len(images) + 1}.{ext}"
                with open(image_filename, "wb") as f:
                    f.write(image.blob)
                images.append(image_filename)

        print(f"[INFO] Page {i}: 텍스트 {len(texts)}개, 표 {len(tables)}개, 이미지 {len(images)}개 추출 완료.")

        # 슬라이드 전체 스냅샷 PNG 생성
        state = export_slide_as_png(state, i)
        slide_image_path = state["slide_image"]
        print(f"[INFO] Page {i}: PNG 생성 완료 -> {slide_image_path}")

        slide_data = SlideData(
            page=i,
            slide_image=slide_image_path,
            texts=texts,
            images=images,
            tables=tables,
        )
        slides_data.append(slide_data)

    print(f"[INFO] 총 {len(prs.slides)} 슬라이드 처리 완료.")
    state["slides"] = slides_data
    return state


# ------------------------------------------------------------
# 모듈 단독 테스트용 (원하면 사용)
# ------------------------------------------------------------
if __name__ == "__main__":
    # 간단 테스트 예시 (직접 실행 시 사용)
    sample_state: State = {
        "pptx_path": "sample.pptx",
        "work_dir": "./work_dir",
        "media_dir": "./work_dir/media",
    }
    os.makedirs(sample_state["media_dir"], exist_ok=True)
    node_parse_ppt(sample_state)
